import streamlit as st
import json
import random
import hashlib
import os
import re
import time
import requests
from datetime import datetime
from contextlib import contextmanager


from streamlit_local_storage import LocalStorage
import fitz
import docx
import pptx
import csv
import io
from concurrent.futures import ThreadPoolExecutor

import pyrebase
import re

try:
    API_KEY = st.secrets["jsonbin"]["api_key"]
    BIN_ID = st.secrets["jsonbin"]["bin_id"]
    BASE_URL = f"https://api.jsonbin.io/v3/b/{BIN_ID}"
    HEADERS = {
        'Content-Type': 'application/json',
        'X-Master-Key': API_KEY
    }

except (KeyError, TypeError):
    # This will show an error on the page if secrets are not set correctly
    st.error("JSONBin credentials not found in st.secrets. Please check your secrets.toml file.")
    st.stop()

config = {
    "apiKey": st.secrets["firebase"]["firebase_apiKey"],
    "authDomain": st.secrets["firebase"]["firebase_authDomain"],
    "databaseURL": st.secrets["firebase"]["firebase_databaseURL"],
    "projectId": st.secrets["firebase"]["firebase_projectId"],
    "storageBucket": st.secrets["firebase"]["firebase_storageBucket"],
}

# --- Initialize Firebase ---
# We use a try-except block to prevent re-initialization on every rerun
try:
    firebase = pyrebase.initialize_app(config)
    db = firebase.database()
except ValueError:
    # If the app is already initialized, just get the instance
    # This is a workaround for Streamlit's execution model
    app = firebase.get_app()
    db = firebase.database(app)

# Configure page
st.set_page_config(
    page_title="Knowledge Quest",
    page_icon="🎓",
    layout="centered",
    initial_sidebar_state="collapsed"
)

localS = LocalStorage()

# Poe API Client
class PoeAPIClient:
    def __init__(self, api_key):
        self.api_key = api_key
        self.base_url = "https://api.poe.com/v1"

    def generate_questions(self, prompt, model="Gemini-3-Flash"):
        """Generate questions using Poe API"""
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }
        data = {
            "model": model,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.7,
            "max_tokens": 20000,
            "stream": False
        }
        try:
            response = requests.post(
                f"{self.base_url}/chat/completions",
                headers=headers,
                json=data,
                timeout=120
            )
            response.raise_for_status()  # Raise an exception for bad status codes (4xx or 5xx)
            result = response.json()
            return result['choices'][0]['message']['content']
        except requests.exceptions.RequestException as e:
            st.error(f"API Request failed: {str(e)}")
            return None

    def generate_tts(self, text, voice="default"):
        """
        Generate TTS audio using Poe API.
        NOTE: This function assumes the Poe API returns an audio URL in an 'attachments'
        field for a chat completion call to a model like 'ElevenLabs-v3'. This is an
        unconventional way to get TTS and may be fragile or incorrect depending on the
        actual Poe API specification.
        """
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }
        data = {
            "model": "ElevenLabs-v3",  # This model name is hypothetical
            "messages": [{"role": "user", "content": text.strip()}],
            "stream": False
        }
        try:
            response = requests.post(
                f"{self.base_url}/chat/completions",
                headers=headers,
                json=data,
                timeout=30
            )
            response.raise_for_status()
            result = response.json()
            message = result.get('choices', [{}])[0].get('message', {})
            if 'attachments' in message:
                for attachment in message['attachments']:
                    if attachment.get('content_type', '').startswith('audio/'):
                        return attachment.get('url')
            # Fallback if the structure is different but URL is in content
            content = message.get('content', '')
            url_match = re.search(r'https?://[^\s]+', content)
            if url_match:
                return url_match.group(0)
            return None
        except requests.exceptions.RequestException as e:
            st.error(f"TTS Generation failed: {e}")
            return None


# --- Improved Horizontal Layout CSS and Context Manager ---
# --- Improved Horizontal Layout CSS and Context Manager ---
# --- Improved Horizontal Layout CSS and Context Manager ---
# --- Improved Horizontal Layout CSS and Context Manager ---
HORIZONTAL_STYLE = """
<style class="hide-element">
    /* Hides the style container itself */
    .element-container:has(.hide-element) {
        display: none;
    }
    /* 
        This is the main selector. It finds a div that has a direct child 
        with class 'horizontal-marker' and applies flexbox styling to it.
    */
    div[data-testid="stVerticalBlock"]:has(> .element-container .horizontal-marker) {
        display: flex;
        flex-direction: row !important;
        align-items: flex-start;
        width: 100%;
        gap: 1rem;
    }
    /* 
        Overrides the default fixed width of Streamlit's container elements 
        within our horizontal block to allow them to shrink-to-fit their content.
    */
    div[data-testid="stVerticalBlock"]:has(> .element-container .horizontal-marker) div {
        width: max-content !important;
    }
    /* Make first and last elements take equal space to center the middle */
    div[data-testid="stVerticalBlock"]:has(> .element-container .horizontal-marker) div:nth-child(2),
    div[data-testid="stVerticalBlock"]:has(> .element-container .horizontal-marker) div:nth-child(4) {
        flex: 1;
    }
    /* Center the middle element and align it to top */
    div[data-testid="stVerticalBlock"]:has(> .element-container .horizontal-marker) div:nth-child(3) {
        flex: 0 0 auto;
        text-align: center;
        margin-top: -5px;
        padding-top: 0 !important;
    }
    /* Align last button to the right */
    div[data-testid="stVerticalBlock"]:has(> .element-container .horizontal-marker) div:nth-child(4) {
        display: flex;
        justify-content: flex-end;
    }
</style>
"""


@contextmanager
def st_horizontal():
    """
    A context manager to layout Streamlit elements horizontally with full width distribution.
    It injects the necessary CSS and a marker span to activate the styling.
    """
    # Inject the CSS for horizontal layout.
    st.markdown(HORIZONTAL_STYLE, unsafe_allow_html=True)
    with st.container():
        # Add a hidden marker element. The CSS selector uses this marker to identify
        # the container that should have its children laid out horizontally.
        st.markdown('<span class="hide-element horizontal-marker"></span>', unsafe_allow_html=True)
        yield


# Custom CSS
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Lato:wght@400;600;700;800;900&display=swap');

/* ... [Rest of your CSS is unchanged and correct] ... */
.main-title{font-family:'Lato',sans-serif;font-size:2.25rem;font-weight:900;text-align:center;margin-bottom:1.5rem;background:linear-gradient(135deg, #4F46E5, #3B82F6);-webkit-background-clip:text;-webkit-text-fill-color:transparent;background-clip:text;}.quiz-container{font-family:'Lato',sans-serif;background:white;border-radius:16px;padding:30px;border:1px solid #D1D5DB;box-shadow:0 8px 25px rgba(0,0,0,0.05);margin:1rem 0;}.question-text{font-size:1.1rem;font-weight:600;margin-bottom:0.75rem;}.progress-container{background-color:#E5E7EB;border-radius:4px;height:8px;margin:1rem 0;}.progress-bar{background:linear-gradient(45deg, #4F46E5, #3B82F6);height:100%;border-radius:4px;transition:width 0.5s ease;}.api-status{background:linear-gradient(45deg, #EBF8FF, #DBEAFE);border:1px solid #3B82F6;border-radius:12px;padding:15px;margin:10px 0;}.demo-questions{background:#F0FDF4;border:1px solid #22C55E;border-radius:8px;padding:15px;margin:10px 0;}.badge{display:inline-block;margin:5px;padding:8px 16px;border-radius:20px;font-size:0.9rem;font-weight:600;}.perfect-badge{background:linear-gradient(45deg, #FBBF24, #F59E0B);color:#78350F;}.high-achiever-badge{background:linear-gradient(45deg, #10B981, #059669);color:white;}.explanation-box{padding:15px;border-radius:12px;border:1px solid #D1D5DB;margin-top:15px;}
</style>
""", unsafe_allow_html=True)


# Initialize session state
def init_session_state():
    defaults = {
        'questions': [], 'original_questions': [], 'current_question_index': 0,
        'questions_completed': 0, 'incorrect_question_ids': set(), 'quiz_started': False,
        'quiz_finished': False, 'user_answers': {}, 'score_history': [],
        'revision_mode': False, 'revision_index': 0, 'audio_urls': {'questions': {}, 'answers': {}},
        'generating_questions': False, 'poe_client': None, 'show_ai_settings': False,
        'audio_generated': False, 'is_redoing_wrong': False, 'quiz_generation_in_progress': False,
        'uploader_key': 0, 'confirm_clear_local': False, 'confirm_clear_cloud': False, 'char_limit': 30000
    }
    for key, default_value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default_value


# Utility functions
def get_poe_api_key():
    api_key = os.getenv("POE_API_KEY")
    if not api_key:
        try:
            api_key = st.secrets.get("POE_API_KEY", "")
        except Exception:
            api_key = ""
    return api_key


def init_poe_client():
    if st.session_state.poe_client is None:
        api_key = get_poe_api_key()
        if api_key:
            st.session_state.poe_client = PoeAPIClient(api_key)
            return True
        return False
    return True


def stable_hash(text):
    return 'q_' + hashlib.md5(text.encode()).hexdigest()[:8]


def strip_markdown_fences(text):
    """
    Extracts content from markdown code fences (```json ... ```) or
    the largest apparent JSON array/object from a string.
    Does NOT validate JSON completeness, just extracts the raw string.
    """
    # 1. Prioritize markdown code fences
    match = re.search(r"```(?:json)?\s*\n(.*?)\n\s*```", text, flags=re.DOTALL)
    if match:
        return match.group(1).strip()

    # 2. Fallback Strategy: Look for a bare JSON array or object
    # We need to find the outermost array or object if no fences are present.
    # This is trickier because the text might contain other prose.

    best_json_candidate = None

    # Try array first
    start_bracket = text.find('[')
    end_bracket = text.rfind(']')
    if start_bracket != -1 and end_bracket > start_bracket:
        candidate = text[start_bracket: end_bracket + 1]
        # Heuristic: does it look like an array of objects?
        if candidate.count('{') > 0 and candidate.count('}') > 0:
            best_json_candidate = candidate

    # Then try object (only if no good array candidate or if object is more prominent)
    start_brace = text.find('{')
    end_brace = text.rfind('}')
    if start_brace != -1 and end_brace > start_brace:
        candidate = text[start_brace: end_brace + 1]
        # If we have an array candidate, and this object is not the *entire* content,
        # we stick with the array. This is a heuristic.
        if best_json_candidate is None or (
                len(candidate) > len(best_json_candidate) and not best_json_candidate.startswith('[')):
            best_json_candidate = candidate

    return best_json_candidate


def parse_partial_json_array(json_string):
    """
    Attempts to parse a JSON string, potentially truncated, to extract as many
    complete JSON objects from an array as possible.
    Returns a list of parsed objects.
    """
    if not json_string:
        return []

    json_string = json_string.strip()

    # If it's a complete, valid JSON array, parse it directly
    try:
        data = json.loads(json_string)
        if isinstance(data, list):
            return data
    except json.JSONDecodeError:
        pass  # It's partial or malformed, proceed to recovery logic

    # If it doesn't start with an array, it's not what we expect for questions.
    # Or if it's a single object, we can try to wrap it, but for a list of questions,
    # we primarily expect an array.
    if not json_string.startswith('['):
        return []

    recovered_objects = []
    balance = 0
    in_string = False
    escape_char = False
    start_obj_index = -1

    # Iterate through the string to find complete top-level objects within the array
    # We need to handle nested structures and strings correctly to avoid misinterpreting '}'
    for i, char in enumerate(json_string):
        if in_string:
            if char == '\\':
                escape_char = not escape_char
            elif char == '"' and not escape_char:
                in_string = False
            else:
                escape_char = False
        elif char == '"':
            in_string = True
            escape_char = False
        elif char == '{':
            if balance == 0:  # Start of a new top-level object in the array
                start_obj_index = i
            balance += 1
        elif char == '}':
            balance -= 1
            # If balance returns to 0 and we had a starting object index,
            # it means we've found a complete top-level object.
            if balance == 0 and start_obj_index != -1:
                try:
                    obj_str = json_string[start_obj_index: i + 1]
                    obj = json.loads(obj_str)
                    recovered_objects.append(obj)
                    start_obj_index = -1  # Reset for the next object
                except json.JSONDecodeError:
                    # This specific object might be malformed despite balanced braces,
                    # or the truncation happened within it. We skip it.
                    pass
    return recovered_objects

def is_valid_json_input(text):
    try:
        parsed = json.loads(strip_markdown_fences(text.strip()))
        return validate_questions_array(parsed)['valid']
    except (json.JSONDecodeError, TypeError):
        return False


def check_input_and_show_ai_settings():
    input_text = st.session_state.get('question_input', '').strip()
    st.session_state.show_ai_settings = bool(input_text and not is_valid_json_input(input_text))


def validate_questions_array(data):
    if not isinstance(data, list) or not data:
        return {'valid': False, 'error': 'Input must be a non-empty JSON array.'}
    for i, q in enumerate(data):
        if not isinstance(q, dict):
            return {'valid': False, 'error': f'Item {i + 1} is not an object.'}
        required = ['question', 'options', 'correct', 'explanation']
        for field in required:
            if field not in q:
                return {'valid': False, 'error': f'Item {i + 1} is missing field: "{field}".'}
        if not (isinstance(q['question'], str) and q['question'].strip()):
            return {'valid': False, 'error': f'Item {i + 1}: "question" must be a non-empty string.'}
        if not (isinstance(q['options'], list) and len(q['options']) >= 2):
            return {'valid': False, 'error': f'Item {i + 1}: "options" must be an array with at least 2 items.'}
        if not (isinstance(q['correct'], int) and 0 <= q['correct'] < len(q['options'])):
            return {'valid': False, 'error': f'Item {i + 1}: "correct" index is invalid.'}
        if 'hint' not in q: q['hint'] = ""
    return {'valid': True}


def generate_ai_prompt(input_text, num_questions):
    return f"""You are a teacher creating educational assessments. You are Usage from Chiikawa, the very cute crazy rabbit character. Let's learn something new!

When questions are asked, you give constructive step by step hints to lead the student to get to the answer, before giving the direct answers but you make sure the student can get to the point at the end. The style of teaching is concise and get to the point, but keep it friendly. When giving compliments and acting like the character, you can use Japanese for non technical related sentence. When you are talking on technical items, please always use English.

You are given the following materials. As images may not be included, you may need to guess what could be related in the materials.

You are a teacher creating educational assessments. Based on the following materials, create {num_questions} multiple-choice questions.

---
{input_text}
---

Based on these materials, do the following:
1) Guess the educational level of the topic (e.g., primary P.2, secondary, tertiary, professional, postgraduate).
2) Create {num_questions} multiple-choice questions whose difficulty is one level harder than the guessed level (e.g., guessed P.2 -> produce P.3-level difficulty or slightly higher). Make them slightly tricky but fair.
3) For each question, write plausible distractors that are GENERALLY INCORRECT (not just wrong relative to this passage). Distractors should represent common misconceptions or confusable alternatives that would be wrong in most contexts.
4) The "explanation" field should be concise and help memorization (shown after correct).
5) The "hint" field must be present (can be short) and should guide reflection after a wrong attempt.
6) Ensure no distractor is a case/spacing variant of the correct answer with similar length.
7) Distractors must be substantively different from the correct answer
8) Use the same language as of the materials given. Only when the materials is about language learning, supplement with English so user get understood everything even user cannot understand the language materials.
9) Note markdowns or codes or mathematical formulas are not rendered on the question (but rendered in options), so the question must be readable in plain text

Important formatting rules:
- Output ONLY pure JSON. No thoughts. No preface. No prose, no markdown, no code fences.
- The JSON must be an array of objects with exactly these keys per item:
question (string), options (array of 3-6 strings), correct (integer index within options), hint (string), explanation (string).
- Do NOT include any extra wrapper objects or metadata. No backticks. No comments.

Example JSON:
[
{{
"question": "What is the capital city of France?",
"options": [
"London",
"Paris", 
"Berlin",
"Madrid"
],
"correct": 1,
"hint": "Think about the most famous city in France.",
"explanation": "Paris is the capital and largest city of France."
}}
]

Provide ONLY the JSON array.
"""

def get_demo_questions():
    return [
        {"question": "What is the capital city of France?", "options": ["London", "Paris", "Berlin", "Madrid"],
         "correct": 1, "hint": "Think about the most famous city in France.",
         "explanation": "Paris is the capital and largest city of France, known for landmarks like the Eiffel Tower."},
        {"question": "Which continent is Brazil located in?",
         "options": ["North America", "South America", "Africa", "Asia"], "correct": 1,
         "hint": "Brazil is the largest country in its continent.",
         "explanation": "Brazil is located in South America and is the continent's largest country by both area and population."},
        {"question": "What is the longest river in the world?",
         "options": ["Amazon River", "Nile River", "Mississippi River", "Yangtze River"], "correct": 1,
         "hint": "This river flows through northeastern Africa.",
         "explanation": "The Nile River is traditionally considered the longest river in the world, flowing through northeastern Africa."}
    ]


# Main quiz functions
def generate_questions_with_ai(input_text, num_questions, model):
    if not st.session_state.poe_client:
        st.error("Poe API client not initialized. Please check your API key.")
        return None

    prompt = generate_ai_prompt(input_text, num_questions)
    with st.spinner(f"🤖 Generating {num_questions} questions with {model}..."):
        response = st.session_state.poe_client.generate_questions(prompt, model)
        if not response:
            st.error("No response received from AI.")
            return None

        # --- FIX: Robust JSON extraction and partial parsing ---
        cleaned_response_str = strip_markdown_fences(response.strip())

        if cleaned_response_str is None:
            st.error("AI response did not contain a recognizable JSON structure.")
            with st.expander("Raw AI Response"): st.text(response)
            return None

        # Use the new function to parse potentially partial JSON
        questions = parse_partial_json_array(cleaned_response_str)

        if not questions:
            st.error("Failed to parse any complete questions from AI response.")
            with st.expander("Cleaned AI Response (potentially partial)"):
                st.text(cleaned_response_str)
            return None

        validation = validate_questions_array(questions)
        if validation['valid']:
            # If the number of parsed questions is less than requested, inform the user.
            if len(questions) < num_questions:
                st.warning(f"💡 Only {len(questions)} out of {num_questions} questions were successfully generated and parsed due to an incomplete AI response. Consider reducing the requested number of questions.")
            return questions
        else:
            st.error(f"Generated questions validation failed for parsed questions: {validation['error']}")
            with st.expander("Parsed (but invalid) Questions"):
                st.json(questions)  # Show the partially parsed questions for debugging
            return None


def generate_audio_for_questions():
    if st.session_state.get('quiz_mode') != 'audio' or not st.session_state.poe_client:
        st.session_state.audio_generated = True
        return

    speak_q = st.session_state.get('speak_question', True)
    speak_a = st.session_state.get('speak_answer', True)
    if not (speak_q or speak_a):
        st.session_state.audio_generated = True
        return

    tasks = []
    for q in st.session_state.questions:
        if speak_q and q['id'] not in st.session_state.audio_urls['questions']:
            tasks.append(('question', q['id'], q['question']))
        if speak_a and q['id'] not in st.session_state.audio_urls['answers']:
            correct_answer = q['options'][q['correct']]
            tasks.append(('answer', q['id'], f"The correct answer is: {correct_answer}"))

    if not tasks:
        st.session_state.audio_generated = True
        return

    progress_bar = st.progress(0, text="🎵 Generating audio...")
    for i, (audio_type, q_id, text) in enumerate(tasks):
        url_key = 'questions' if audio_type == 'question' else 'answers'
        audio_url = st.session_state.poe_client.generate_tts(text)
        if audio_url:
            st.session_state.audio_urls[url_key][q_id] = audio_url
        progress_bar.progress((i + 1) / len(tasks), text=f"🎵 Generating audio... ({i + 1}/{len(tasks)})")

    progress_bar.empty()
    st.session_state.audio_generated = True


def start_quiz():
    input_text = st.session_state.get('question_input', '').strip()

    # --- MODIFICATION START ---
    # Define a character limit for the text to be processed by the AI.
    CHAR_LIMIT = st.session_state.get('char_limit', 20000)

    # Check if the input text exceeds the character limit.
    if len(input_text) > CHAR_LIMIT:
        # Show a non-blocking "pop-up" message to the user.
        st.toast(
            f"Input text exceeded {CHAR_LIMIT:,} characters and was truncated.",
            icon="⚠️"
        )
        # Truncate the text before it's used for JSON parsing or AI generation.
        input_text = input_text[:CHAR_LIMIT]
    # --- MODIFICATION END ---

    print("input_text")
    if not input_text:
        st.error("Please paste some material or a question set to begin.")
        return

    questions = None
    try:
        # Try to parse as JSON first
        stripped_input = strip_markdown_fences(input_text)
        if stripped_input is not None:
            parsed = json.loads(stripped_input)
            validation = validate_questions_array(parsed)
            if validation['valid']:
                questions = parsed
            else:
                # If it's not valid JSON, but was intended as such, show error
                if input_text.strip().startswith('['):
                    st.error(f"JSON format error: {validation['error']}")
                    return
    except json.JSONDecodeError:
        pass  # Not JSON, will proceed to AI generation

    # If not valid JSON, generate with AI
    if questions is None:
        if not init_poe_client():
            st.error("❌ Poe API key not found. Cannot generate questions.")
            st.info("💡 Using demo geography questions instead.")
            questions = get_demo_questions()
        else:
            num_q = st.session_state.get('num_questions', 3)
            model = st.session_state.get('llm_model', 'GPT-5-mini')
            # The (potentially truncated) 'input_text' is used here.
            questions = generate_questions_with_ai(input_text, num_q, model)
            if not questions:
                st.error("Failed to generate questions. Using demo questions instead.")
                questions = get_demo_questions()

    if questions:
        setup_quiz_with_questions(questions)


def setup_quiz_with_questions(questions_data):
    questions = json.loads(json.dumps(questions_data))  # Deep copy
    for q in questions:
        q['id'] = stable_hash(f"{q['question']}|{'|'.join(map(str, q['options']))}|{q['correct']}")

    if not st.session_state.get('is_redoing_wrong', False):
        random.shuffle(questions)
        st.session_state.original_questions = json.loads(json.dumps(questions))

    st.session_state.questions = questions
    st.session_state.quiz_started = True
    st.session_state.quiz_finished = False
    st.session_state.current_question_index = 0
    st.session_state.questions_completed = 0
    st.session_state.user_answers = {}
    st.session_state.audio_generated = False

    if st.session_state.get('quiz_mode') == 'audio':
        generate_audio_for_questions()
    else:
        st.session_state.audio_generated = True
    st.rerun()


def play_audio(key):
    st.session_state.audio_to_play = key


def render_quiz_question():
    if not st.session_state.questions or not st.session_state.audio_generated:
        st.info("Please wait, preparing quiz...")
        st.rerun()

    idx = st.session_state.current_question_index
    if idx >= len(st.session_state.questions):
        show_quiz_summary()
        return

    question = st.session_state.questions[idx]
    total_q = len(st.session_state.questions)
    progress = st.session_state.questions_completed / total_q

    st.markdown(
        f'<div class="progress-container"><div class="progress-bar" style="width: {progress * 100}%"></div></div>',
        unsafe_allow_html=True)

    # Navigation and question counter
    with st_horizontal():
        # "Back" button
        if st.button("← Back", disabled=idx == 0):
            go_back()

        # Question counter
        st.markdown(
            f"<div style='text-align: center; font-weight: 600; padding-top: 0.5rem;'>Question {idx + 1} of {total_q}</div>",
            unsafe_allow_html=True)

        # "Next" or "Finish" button logic
        is_last_question = (idx == total_q - 1)
        answered = question['id'] in st.session_state.user_answers

        if is_last_question:
            if st.button("🏁 Finish Quiz", type="primary", disabled=not answered):
                finish_quiz()
        else:
            if st.button("Next →", disabled=not answered):
                go_next()

    # Question and audio
    st.markdown(f"<div class='question-text'>{question['question']}</div>", unsafe_allow_html=True)
    if st.session_state.get('quiz_mode') == 'audio':
        render_audio_controls(question)

    # Answer options
    render_answer_options(question)

    # Explanation if answered
    if answered:
        show_answer_result(question)


def render_audio_controls(question):
    # --- FIX: Reliable audio playback ---
    # Autoplay is unreliable; use explicit buttons.
    audio_placeholder = st.empty()
    if 'audio_to_play' in st.session_state and st.session_state.audio_to_play:
        audio_placeholder.audio(st.session_state.audio_to_play, format="audio/mp3", autoplay=True)
        # Small delay to allow audio to start, then clear the state to prevent re-playing on rerun
        time.sleep(0.1)
        st.session_state.audio_to_play = None

    c1, c2 = st.columns(2)
    q_audio_url = st.session_state.audio_urls['questions'].get(question['id'])
    if st.session_state.get('speak_question') and q_audio_url:
        if c1.button("🔊 Play Question"):
            play_audio(q_audio_url)

    answered = question['id'] in st.session_state.user_answers
    a_audio_url = st.session_state.audio_urls['answers'].get(question['id'])
    if answered and st.session_state.get('speak_answer') and a_audio_url:
        if c2.button("🔊 Play Answer"):
            play_audio(a_audio_url)


def render_answer_options(question):
    q_id = question['id']
    # Shuffle options once per question and store in the question object itself
    if 'shuffled_options' not in question:
        original_options = list(enumerate(question['options']))
        random.shuffle(original_options)
        st.session_state.questions[st.session_state.current_question_index]['shuffled_options'] = original_options

    shuffled_options = question['shuffled_options']
    answered = q_id in st.session_state.user_answers
    user_selection = st.session_state.user_answers.get(q_id, {}).get('selected')

    for original_idx, option_text in shuffled_options:
        is_correct = (original_idx == question['correct'])
        is_selected = (original_idx == user_selection)

        button_key = f"option_{q_id}_{original_idx}"
        if st.button(option_text, key=button_key, disabled=answered, use_container_width=True):
            handle_answer_selection(q_id, original_idx)
            st.rerun()

        # --- FIX: Clearer visual feedback after answering ---
        if answered:
            if is_correct:
                st.success(f"✅ Correct Answer: {option_text}")
            elif is_selected:
                st.error(f"❌ Your Choice: {option_text}")


def handle_answer_selection(q_id, selected_option_idx):
    idx = st.session_state.current_question_index
    question = st.session_state.questions[idx]
    is_correct = (selected_option_idx == question['correct'])

    if q_id not in st.session_state.user_answers:
        st.session_state.user_answers[q_id] = {
            'selected': selected_option_idx,
            'is_correct': is_correct,
            'first_try': True
        }
        if is_correct:
            st.session_state.questions_completed += 1
        else:
            # --- FIX: Robust incorrect question tracking ---
            st.session_state.incorrect_question_ids.add(q_id)
            st.session_state.questions_completed += 1
    # This function is now simpler; visual feedback is handled in render_answer_options


def show_answer_result(question):
    answer_record = st.session_state.user_answers.get(question['id'])
    if answer_record and answer_record['is_correct']:
        if question.get('explanation'):
            st.markdown(
                f'<div class="explanation-box"><strong>Explanation:</strong><br>{question["explanation"]}</div>',
                unsafe_allow_html=True)
    elif answer_record:  # Incorrect answer
        if question.get('hint'):
            st.warning(f"💡 **Hint:** {question['hint']}")


def go_back():
    if st.session_state.current_question_index > 0:
        st.session_state.current_question_index -= 1
        st.rerun()


def go_next():
    if st.session_state.current_question_index < len(st.session_state.questions) - 1:
        st.session_state.current_question_index += 1
        st.rerun()


def finish_quiz():
    st.session_state.quiz_finished = True
    st.rerun()


def show_quiz_summary():
    st.markdown("<div class='main-title'>✨ Quiz Complete! ✨</div>", unsafe_allow_html=True)
    total_q = len(st.session_state.questions)
    correct_answers = len([ans for ans in st.session_state.user_answers.values() if ans['is_correct']])
    score = round((correct_answers / total_q) * 100) if total_q > 0 else 0

    emoji, message = ("🎉", "Perfect score!") if score == 100 else \
        ("🌟", "Excellent work!") if score >= 80 else \
            ("👍", "Well done!") if score >= 60 else \
                ("💪", "Good effort!")

    st.markdown(f"""
    <div style='text-align: center; margin: 2rem 0;'>
        <div style='font-size: 2rem;'>{emoji}</div>
        <div style='font-size: 1.5rem; font-weight: 700; margin: 10px 0;'>{message}</div>
        <div style='font-size: 1.2rem;'>Score: {correct_answers} / {total_q} ({score}%)</div>
    </div>
    """, unsafe_allow_html=True)

    c1, c2, c3, c4 = st.columns(4)
    c1.button("🔄 Retry All", on_click=retry_quiz, use_container_width=True)
    c2.button("❌ Redo Wrong", on_click=redo_incorrect_questions, use_container_width=True,
              disabled=not st.session_state.incorrect_question_ids)
    c3.button("📚 Review", on_click=start_revision_mode, use_container_width=True,
              disabled=not st.session_state.incorrect_question_ids)
    c4.button("🆕 New Quiz", on_click=clear_quiz, use_container_width=True)


def reset_quiz_state():
    """Helper to reset common state variables for a new quiz attempt."""
    st.session_state.current_question_index = 0
    st.session_state.questions_completed = 0
    st.session_state.user_answers = {}
    st.session_state.quiz_finished = False
    st.session_state.is_redoing_wrong = False  # --- FIX: Ensure this is always reset ---
    st.session_state.incorrect_question_ids = set()


def retry_quiz():
    reset_quiz_state()
    # Reshuffle the original set of questions
    questions = json.loads(json.dumps(st.session_state.original_questions))
    random.shuffle(questions)
    st.session_state.questions = questions
    # Clear shuffled options from previous attempts
    for q in st.session_state.questions:
        q.pop('shuffled_options', None)


def redo_incorrect_questions():
    if not st.session_state.incorrect_question_ids:
        st.warning("No incorrect questions to redo!")
        return

    # Filter original questions to get only the incorrect ones
    incorrect_q_ids = st.session_state.incorrect_question_ids
    wrong_questions = [q for q in st.session_state.original_questions if q['id'] in incorrect_q_ids]

    reset_quiz_state()
    st.session_state.is_redoing_wrong = True
    st.session_state.questions = wrong_questions
    # Clear shuffled options from previous attempts
    for q in st.session_state.questions:
        q.pop('shuffled_options', None)


def start_revision_mode():
    if not st.session_state.incorrect_question_ids:
        st.warning("No incorrect questions to review!")
        return
    st.session_state.revision_mode = True
    st.session_state.revision_index = 0


def clear_quiz():
    # Preserve API client and key
    client = st.session_state.get('poe_client')
    # Clear all other session state keys
    keys_to_clear = [k for k in st.session_state.keys() if k != 'poe_client']
    for key in keys_to_clear:
        del st.session_state[key]
    init_session_state()
    st.session_state.poe_client = client  # Restore client


def render_revision_mode():
    st.markdown("<div class='main-title'>📚 Revision Mode</div>", unsafe_allow_html=True)

    incorrect_q_ids = list(st.session_state.incorrect_question_ids)
    if not incorrect_q_ids:
        st.session_state.revision_mode = False
        st.rerun()
        return

    idx = st.session_state.revision_index
    q_id = incorrect_q_ids[idx]
    # Find the question from the original list
    question = next((q for q in st.session_state.original_questions if q['id'] == q_id), None)

    if not question:
        st.error("Error: Could not find question to review.")
        st.session_state.revision_mode = False
        return

    # --- MODIFIED: Use st_horizontal for revision navigation with full width ---
    with st_horizontal():
        if st.button("← Previous", disabled=idx == 0):
            st.session_state.revision_index -= 1
            st.rerun()

        st.markdown(
            f"<div style='text-align: center; font-weight: 600; padding-top: 0.5rem;'>Reviewing {idx + 1} of {len(incorrect_q_ids)}</div>",
            unsafe_allow_html=True)

        if st.button("Next →", disabled=idx >= len(incorrect_q_ids) - 1):
            st.session_state.revision_index += 1
            st.rerun()

    st.markdown(f"<div class='question-text'>{question['question']}</div>", unsafe_allow_html=True)
    st.markdown("**Correct Answer:**")
    st.success(f"✅ {question['options'][question['correct']]}")
    if question.get('explanation'):
        st.markdown(f'<div class="explanation-box"><strong>Explanation:</strong><br>{question["explanation"]}</div>',
                    unsafe_allow_html=True)

    if st.button("✅ Finish Revision", use_container_width=True, type="primary"):
        st.session_state.revision_mode = False
        st.rerun()


# --- TEXT COLLECTOR HELPER FUNCTIONS ---

def tc_initialize_state():
    """
    Initializes session state from browser's local storage.
    This function is designed to handle the asynchronous nature of fetching data
    from the browser, which often requires more than one script run on initial load.
    """
    # 1. Fetch the item from local storage. It might be None on the first run after a refresh.
    persisted_json = localS.getItem("all_texts")

    if 'all_texts' not in st.session_state:
        st.session_state.all_texts = json.loads(persisted_json) if persisted_json else {}
        st.session_state.processed_files = set(st.session_state.all_texts.keys())

    elif not st.session_state.all_texts and persisted_json:
        st.session_state.all_texts = json.loads(persisted_json)
        st.session_state.processed_files = set(st.session_state.all_texts.keys())



def tc_save_data():
    """Saves the collected texts to local storage."""
    localS.setItem("all_texts", json.dumps(st.session_state.all_texts))


def tc_extract_text_from_file(uploaded_file):
    """Extracts text content from a supported file type."""
    name, extension = os.path.splitext(uploaded_file.name)
    extension = extension.lower()
    if extension == ".pdf":
        with fitz.open(stream=uploaded_file.getvalue(), filetype="pdf") as doc:
            return "".join(page.get_text() for page in doc)
    elif extension == ".docx":
        doc = docx.Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])
    elif extension == ".pptx":
        """
        Extracts text from a .pptx file.
        This will extract text from all shapes on each slide.
        """
        try:
            prs = pptx.Presentation(uploaded_file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            return "\n".join(text_runs)
        except Exception as e:
            st.error(f"Error processing PPTX file: {e}")
            return None
    elif extension == ".csv":
        """
        Extracts text from a .csv file.
        Converts rows and columns into a single block of text.
        """
        try:
            # Read the file content and decode it
            content = uploaded_file.getvalue().decode("utf-8")
            # Use io.StringIO to allow the csv module to read a string
            file_like_object = io.StringIO(content)
            reader = csv.reader(file_like_object)
            # Join cells with a space, and rows with a newline
            return "\n".join([" ".join(row) for row in reader])
        except Exception as e:
            st.error(f"Error processing CSV file: {e}")
            return None

    elif extension == ".json":
        """
        Extracts and formats text from a .json file.
        This will "pretty-print" the JSON for readability.
        """
        try:
            # Read the file content and decode it
            content = uploaded_file.getvalue().decode("utf-8")
            # Parse the JSON data
            data = json.loads(content)
            # Convert it back to a formatted string with indentation
            return json.dumps(data, indent=4)
        except json.JSONDecodeError:
            st.error("Error: The uploaded JSON file is not correctly formatted.")
            return None
        except Exception as e:
            st.error(f"Error processing JSON file: {e}")
            return None
    elif extension in [".txt", ".md"]:
        return uploaded_file.getvalue().decode("utf-8")
    else:
        st.error(f"Unsupported file type: {extension}")
        return None


@st.cache_data(ttl=60) # Cache for 1 minute to avoid hitting API rate limits
# def get_all_cloud_data():
#     """Reads the entire database (one JSON file) from JSONBin."""
#     try:
#         response = requests.get(f"{BASE_URL}/latest", headers=HEADERS)
#         # Raise an exception for bad status codes (4xx or 5xx)
#         response.raise_for_status()
#         # The actual data is nested under the "record" key
#         return response.json().get("record", {})
#     except requests.exceptions.HTTPError as e:
#         # A 404 error is expected if the bin is new or empty. Treat as empty data.
#         if e.response.status_code == 404:
#             return {}
#         st.error(f"Failed to read from JSONBin: {e}")
#         return None
#     except Exception as e:
#         st.error(f"An unexpected error occurred while fetching data: {e}")
#         return None
def get_all_cloud_data():
    """
    Fetches the entire dataset from Firebase Realtime Database.
    The data is expected to be stored under a main 'users' node.
    """
    try:
        # db.child("users") points to the main data node.
        # .get().val() fetches the value.
        response = db.child("users").get().val()

        # If the database is empty or the node doesn't exist, it returns None
        if response is None:
            return {}  # Return an empty dict to match the app's expectation

        return response

    except Exception as e:
        st.error(f"Failed to load cloud data from Firebase: {e}")
        # Return None to indicate a failure in the connection or rules
        return None


# def save_all_cloud_data(data):
#     """Saves the entire database (one JSON file) back to JSONBin."""
#     try:
#         response = requests.put(BASE_URL, json=data, headers=HEADERS)
#         response.raise_for_status()
#         # Clear the cache after a successful write to ensure the next read gets the fresh data
#         st.cache_data.clear()
#         return True
#     except Exception as e:
#         st.error(f"Failed to save data to JSONBin: {e}")
#         return False
def save_all_cloud_data(data):
    """
    Saves/Updates data to Firebase Realtime Database using the .update() method.
    This is highly efficient as it only sends the changed data, not the whole object.
    """
    try:
        # .update() is the magic here. It merges the provided data with the
        # existing data in the cloud. If you upload 5 new files, it only
        # sends those 5 files, solving the "payload too big" 403 error.
        db.child("users").update(data)
        return True

    except Exception as e:
        st.error(f"Failed to save data to Firebase: {e}")
        # This might happen if security rules are wrong or network is down.
        return False

def sanitize_firebase_key(key: str) -> str:
    """Replaces Firebase-invalid characters ('.', '$', '#', '[', ']', '/') with underscores."""
    if not isinstance(key, str):
        return "" # Return empty string if key is not a string
    return re.sub(r'[.#$\[\]/]', '_', key)

# --- MAIN TEXT COLLECTOR PAGE RENDER FUNCTION ---

def render_text_collector_page():
    """
    Renders the full UI, with fixes for delete button layout, functionality,
    and a new "Clear Cloud Data" feature. This version uses st.rerun() for all
    state updates, eliminating the need for hard page reloads.
    """
    tc_initialize_state()

    if st.button("⬅️ Back to Knowledge Quest"):
        st.session_state.page = 'main'
        st.rerun()

    st.title("📚 Text Collector")
    st.markdown(
        "Add documents from your device or paste text. If you enter a Resources ID, your sources will be automatically saved to the cloud."
    )

    # --- Cloud Storage Section ---
    st.subheader("☁️ Cloud Storage")
    user_id = st.text_input("Enter a Resources ID to save/load your sources online.", key="user_id", placeholder="e.g., alex123")
    if user_id:
        user_id = sanitize_firebase_key(user_id)

    if st.button("🔄 Load My Sources from Cloud", disabled=not user_id):
        if user_id:
            all_cloud_data = get_all_cloud_data()
            if all_cloud_data is not None:
                user_sources = all_cloud_data.get(user_id, {})
                if user_sources:
                    st.session_state.all_texts.update(user_sources)
                    st.session_state.processed_files.update(user_sources.keys())
                    tc_save_data()
                    st.success(f"Loaded {len(user_sources)} source(s) from the cloud for '{user_id}'.")
                    st.rerun()
                else:
                    st.info("No sources found in the cloud for this User ID.")
        else:
            st.warning("Please enter a Resources ID to load data from the cloud.")

    st.divider()

    col1, col2 = st.columns([1, 1.3])

    # --- Column 1: Add New Source ---
    with col1:
        st.header("1. Add New Source")
        tab1, tab2 = st.tabs(["📄 Upload File(s)", "📋 Paste Text"])

        # This callback function for file uploads is complex but appears logically sound.
        # No major changes are needed here based on the reported issues.
        def handle_file_upload(uploader_key):
            if uploader_key in st.session_state and st.session_state[uploader_key]:
                uploaded_files = st.session_state[uploader_key]
                files_to_process = [
                    file for file in uploaded_files
                    if file.name not in st.session_state.processed_files
                ]

                if not files_to_process:
                    st.info("No new files to process.")
                    st.session_state['uploader_key'] += 1
                    return

                success_results = {}
                failed_files = {}
                processed_file_names = set()

                with st.spinner(f"Processing {len(files_to_process)} new file(s)..."):
                    with ThreadPoolExecutor(max_workers=5) as executor:
                        future_to_file = {executor.submit(tc_extract_text_from_file, file): file for file in files_to_process}
                        for future in future_to_file:
                            file = future_to_file[future]
                            try:
                                content = future.result()
                                if content is not None:
                                    sanitized_name = sanitize_firebase_key(file.name)
                                    success_results[sanitized_name] = content
                            except Exception as e:
                                failed_files[file.name] = str(e)
                            finally:
                                processed_file_names.add(file.name)

                if success_results:
                    st.session_state.all_texts.update(success_results)
                    st.session_state.processed_files.update(processed_file_names)
                    tc_save_data()
                    st.success(f"Successfully added {len(success_results)} file(s) locally.")

                    if user_id:
                        with st.spinner(f"Saving {len(success_results)} file(s) to the cloud..."):
                            all_cloud_data = get_all_cloud_data()
                            if all_cloud_data is not None:
                                if user_id not in all_cloud_data:
                                    all_cloud_data[user_id] = {}
                                all_cloud_data[user_id].update(success_results)
                                if save_all_cloud_data(all_cloud_data):
                                    st.success(f"Successfully saved {len(success_results)} new file(s) to the cloud.")
                                else:
                                    st.error("Failed to save to the cloud due to size limits or network issues.")

                if failed_files:
                    for name, error_msg in failed_files.items():
                        st.error(f"Error processing '{name}': {error_msg}")

                st.session_state['uploader_key'] += 1
                st.rerun()

        with tab1:
            current_uploader_key = f"file_uploader_{st.session_state.get('uploader_key', 0)}"
            st.file_uploader("Upload Documents (PDF, DOCX, etc.)",
                             type=["pdf", "docx", "pptx", "txt", "csv", "json", "md"],
                             accept_multiple_files=True,
                             key=current_uploader_key,
                             on_change=handle_file_upload,
                             kwargs={'uploader_key': current_uploader_key})

        with tab2:
            with st.form("paste_form", clear_on_submit=True):
                source_name = st.text_input("Source Name", placeholder="e.g., 'My Chapter Notes'")
                pasted_text = st.text_area("Paste Text", height=150)
                if st.form_submit_button("Save Pasted Text"):
                    sanitized_name = sanitize_firebase_key(source_name) if source_name else ""
                    if sanitized_name and pasted_text:
                        if sanitized_name in st.session_state.all_texts:
                            st.warning("A source with this name already exists.")
                        else:
                            st.session_state.all_texts[sanitized_name] = pasted_text
                            tc_save_data()
                            st.success(f"Added locally: **{sanitized_name}**")

                            if user_id:
                                with st.spinner("Saving to cloud..."):
                                    all_cloud_data = get_all_cloud_data()
                                    if all_cloud_data is not None:
                                        # --- CORRECTED LOGIC ---
                                        # Ensure user's data object exists, then add/update the single new item.
                                        # This prevents overwriting all other cloud data.
                                        if user_id not in all_cloud_data:
                                            all_cloud_data[user_id] = {}
                                        all_cloud_data[user_id][sanitized_name] = pasted_text

                                        if save_all_cloud_data(all_cloud_data):
                                            st.success(f"Successfully saved '{sanitized_name}' to the cloud.")
                                        else:
                                            st.error(f"Failed to save '{sanitized_name}' to the cloud.")
                            st.rerun()
                    else:
                        st.warning("Please provide both a unique source name and text content.")

        st.markdown("---")

        # --- "Clear All" Buttons with Confirmation ---
        bcol1, bcol2 = st.columns(2)
        with bcol1:
            if st.session_state.get('confirm_clear_local', False):
                st.warning("**Are you sure?** This will delete all local data from your browser.")
                c1, c2 = st.columns(2)
                if c1.button("✅ Yes, Clear Local", use_container_width=True, key="confirm_local_yes"):
                    st.session_state.all_texts = {}
                    st.session_state.processed_files = set()
                    tc_save_data()
                    st.session_state.confirm_clear_local = False
                    st.success("Local data has been cleared.")
                    st.components.v1.html("<script>window.location.reload();</script>", height=0, width=0)
                    # st.rerun() # CORRECT: Use st.rerun() instead of a hard reload.
                if c2.button("❌ No, Cancel", use_container_width=True, key="confirm_local_no"):
                    st.session_state.confirm_clear_local = False
                    st.rerun()
            else:
                if st.button("🗑️ Clear All Local Data", use_container_width=True):
                    st.session_state.confirm_clear_local = True
                    st.rerun()

        with bcol2:
            if st.session_state.get('confirm_clear_cloud', False):
                st.warning(f"**Are you sure?** This will permanently delete all cloud data for '{user_id}'.")
                c1, c2 = st.columns(2)
                if c1.button("✅ Yes, Clear Cloud", use_container_width=True, key="confirm_cloud_yes"):
                    with st.spinner(f"Clearing all cloud data for '{user_id}'..."):
                        all_cloud_data = get_all_cloud_data()
                        if all_cloud_data is not None and user_id in all_cloud_data:
                            all_cloud_data[user_id] = {}
                            if save_all_cloud_data(all_cloud_data):
                                # Also clear local data for consistency
                                st.session_state.all_texts = {}
                                st.session_state.processed_files = set()
                                tc_save_data()
                                st.success(f"Successfully cleared all cloud and local data for '{user_id}'.")
                            else:
                                st.error("Failed to clear cloud data. Local data remains untouched.")
                        else:
                            st.warning("No cloud data found for this user to clear.")
                    st.session_state.confirm_clear_cloud = False
                    st.components.v1.html("<script>window.location.reload();</script>", height=0, width=0)
                    # st.rerun() # CORRECT: Rerun to refresh the UI.
                if c2.button("❌ No, Cancel", use_container_width=True, key="confirm_cloud_no"):
                    st.session_state.confirm_clear_cloud = False
                    st.rerun()
            else:
                if st.button("☁️ Clear My Cloud Data", use_container_width=True, disabled=not user_id):
                    st.session_state.confirm_clear_cloud = True
                    st.rerun()

    # --- Column 2: Combine and Use Text ---
    with col2:
        st.header("2. Combine & Use Text")
        if not st.session_state.all_texts:
            st.info("No sources saved. Upload a file or load from the cloud to begin.")
        else:
            all_doc_names = sorted(list(st.session_state.all_texts.keys()))
            selected_docs = st.multiselect("Choose sources to combine:", options=all_doc_names, key="doc_multiselect")

            # Sync selected sources to cloud (This logic was okay)
            if selected_docs and user_id:
                if st.button(f"⬆️ Sync {len(selected_docs)} Selected Source(s) to Cloud", use_container_width=True):
                    all_cloud_data = get_all_cloud_data()
                    if all_cloud_data is not None:
                        if user_id not in all_cloud_data: all_cloud_data[user_id] = {}
                        for doc_name in selected_docs:
                            all_cloud_data[user_id][doc_name] = st.session_state.all_texts[doc_name]
                        if save_all_cloud_data(all_cloud_data):
                            st.success(f"Successfully synced {len(selected_docs)} source(s) to the cloud.")
                        else:
                            st.error("Failed to sync sources to the cloud.")

            if selected_docs:
                content_blocks = [f"--- Content of: {doc} ---\n\n{st.session_state.all_texts[doc]}" for doc in selected_docs]
                appended_text = "\n\n".join(content_blocks)
                st.text_area(f"Combined Content ({len(selected_docs)} Sources)", appended_text, height=250, key="combined_text_area")

                def use_text_for_quiz_and_switch_view():
                    st.session_state.text_to_inject = st.session_state.combined_text_area
                    st.session_state.page = 'main'
                st.button("🚀 Use this Text for Quiz & Return Home", type="primary", on_click=use_text_for_quiz_and_switch_view, use_container_width=True)

            st.markdown("---")
            st.subheader("Manage Saved Sources")

            # Use a copy of the keys to prevent errors while iterating and deleting
            for doc_name in list(st.session_state.all_texts.keys()):
                c1, c2, c3 = st.columns([0.6, 0.2, 0.2])
                c1.text(doc_name)

                # --- CORRECTED INDIVIDUAL DELETE LOGIC ---
                if c2.button("🗑️ Local", key=f"delete_local_{doc_name}", help="Delete from this browser's storage", use_container_width=True):
                    if doc_name in st.session_state.all_texts:
                        del st.session_state.all_texts[doc_name]
                    st.session_state.processed_files.discard(doc_name)
                    tc_save_data()
                    st.toast(f"Deleted '{doc_name}' from local storage.")
                    st.components.v1.html("<script>window.location.reload();</script>", height=0, width=0)
                    # st.rerun()

                if c3.button("☁️ Cloud", key=f"delete_cloud_{doc_name}", help="Delete from your online cloud storage", use_container_width=True, disabled=not user_id):
                    with st.spinner(f"Deleting '{doc_name}' from cloud..."):
                        all_cloud_data = get_all_cloud_data()
                        if all_cloud_data is not None and user_id in all_cloud_data and doc_name in all_cloud_data[user_id]:
                            del all_cloud_data[user_id][doc_name]
                            if save_all_cloud_data(all_cloud_data):
                                st.success(f"Deleted '{doc_name}' from the cloud.")
                                # Also delete locally for consistency
                                if doc_name in st.session_state.all_texts:
                                    del st.session_state.all_texts[doc_name]
                                    st.session_state.processed_files.discard(doc_name)
                                    tc_save_data()
                            else:
                                st.error(f"Failed to delete '{doc_name}' from the cloud. It remains locally.")
                        else:
                            st.warning(f"'{doc_name}' not found in cloud. Removing it locally.")
                            if doc_name in st.session_state.all_texts:
                                del st.session_state.all_texts[doc_name]
                                st.session_state.processed_files.discard(doc_name)
                                tc_save_data()
                    st.components.v1.html("<script>window.location.reload();</script>", height=0, width=0)
                    # st.rerun()

# Main app
def main():
    init_session_state()
    init_poe_client()

    # We default to the 'main' quiz page.
    if 'page' not in st.session_state:
        st.session_state.page = 'main'

    # If the state is 'text_collector', render that page and stop.
    if st.session_state.page == 'text_collector':
        render_text_collector_page()
        return  # This stops the rest of the main function from running

    # --- Main Quiz App Logic (runs if page is not 'text_collector') ---

    # Inject text from the collector if it exists
    if 'text_to_inject' in st.session_state and st.session_state.text_to_inject:
        st.session_state.question_input = st.session_state.text_to_inject
        del st.session_state.text_to_inject

    if st.session_state.get('revision_mode'):
        render_revision_mode()
    elif st.session_state.get('quiz_finished'):
        show_quiz_summary()
    elif st.session_state.get('quiz_started'):
        render_quiz_question()
    else:
        def set_quiz_generation_status():
            st.session_state.quiz_generation_in_progress = True
        def reset_quiz_generation_status():
            st.session_state.quiz_generation_in_progress = False
        # Setup page
        st.markdown("<div class='main-title'>✨ Knowledge Quest ✨</div>", unsafe_allow_html=True)

        # --- FIX: Conditional AI Settings UI ---
        st.text_area(
            "Paste your materials or pre-formatted JSON to begin:",
            placeholder="Paste any text for AI-generated questions, or paste a valid JSON array of questions.",
            height=200, key="question_input", on_change=check_input_and_show_ai_settings
        )
        char_count = len(st.session_state.get("question_input", ""))
        MAX_CHAR_LIMIT = st.session_state.get("char_limit", 20000)
        # Display the character counter. st.caption is ideal for small helper text.


        # Display a warning if the character count exceeds the limit.
        if char_count > MAX_CHAR_LIMIT:
            st.caption(f"{char_count} / {MAX_CHAR_LIMIT} characters")
            st.warning(
                f"Warning: Your text exceeds the {MAX_CHAR_LIMIT} character limit. "
                f"First {MAX_CHAR_LIMIT} characters is used for processing."
            )


        def go_to_text_collector():
            st.session_state.page = 'text_collector'

        st.button("Or, Extract Text from Documents", icon="📚", on_click=go_to_text_collector)

        # show_ai_panel = st.session_state.get('show_ai_settings', False)
        # if show_ai_panel:
        with st.expander("🤖 AI Generation Settings", expanded=True):
            c1, c2 = st.columns(2)
            c1.selectbox("AI Model:", ["Gemini-3-Flash",
                                       'Grok-4.1-Fast-Reasoning',
                                       # "Gemini-2.5-Pro",
                                       # "GPT-5",
                                       "GPT-5-mini"],
                         key="llm_model",
                         on_change=reset_quiz_generation_status)
            c2.number_input("Number of Questions:", min_value=1, max_value=20, value=3, key="num_questions", on_change=reset_quiz_generation_status)


        st.session_state.quiz_mode = 'Silent Mode'
        # with st.expander("⚙️ Quiz Settings"):
        #     st.selectbox("Quiz Mode:", ["silent", "audio"],
        #                  format_func=lambda x: "Silent Mode" if x == "silent" else "Audio Mode", key="quiz_mode")
        #     if st.session_state.quiz_mode == 'audio':
        #         st.checkbox("Speak question audio", value=True, key="speak_question")
        #         st.checkbox("Speak correct answer audio", value=True, key="speak_answer")

        if st.button("🚀 Start Quiz",
                     type="primary",
                     use_container_width=True,
                     on_click=set_quiz_generation_status,
                     disabled=st.session_state.get('quiz_generation_in_progress', False)):
            start_quiz()


if __name__ == "__main__":
    main()